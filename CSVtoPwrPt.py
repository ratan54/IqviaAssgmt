import numpy as np
import matplotlib.pyplot as mlt
from pptx import Presentation

from pptx.util import Inches

data_file=np.loadtxt('data_file.txt',delimiter=',')

time= data_file[:,0]

time= time -time[0]


sensors= data_file[:,1:5]


avg= np.mean(sensors,0)

avg1= np.mean(sensors,1)

print(avg)

print(avg1)

my_data= np.vstack((time,sensors.T,avg1))
np.savetxt('tempexpo1.csv',my_data,delimiter=',')

my_data= my_data.T

np.savetxt('export1.csv',my_data,delimiter=',')

mlt.plot(time/60,sensors[:,0],'ro')
mlt.plot(time/60,avg1,'b.')


mlt.legend(['Sensor 1','Average'],loc='best')

mlt.xlabel('Time(min)')
mlt.ylabel('Values')
mlt.savefig('fig_1')

mlt.plot(time/60,sensors[:,1],'ro')
mlt.plot(time/60,avg1,'b.')
mlt.legend(['Sensor 2','Average'],loc='best')
mlt.xlabel('Time(min)')
mlt.ylabel('Values')
mlt.savefig('fig_2')


img_path='fig_1.png'
img_path2='fig_2.png'
prs = Presentation()

TSlide=prs.slide_layouts[0]                             # Choosing Title Slide layout
slide=prs.slides.add_slide(TSlide)                      # Adding a slide
title=slide.shapes.title                                # Assigning a title
subtitle=slide.placeholders[1]                          # Adding subtitle placeholder
title.text="Welcome and Hello World!!"                  # Title
subtitle.text="Lets get started"                        # Subtitle

#Adding images

blank_slide_layout = prs.slide_layouts[6]               # Choosing a blank slide layout
slide = prs.slides.add_slide(blank_slide_layout)        # Adding a blank slide
left = top = Inches(1)
pic = slide.shapes.add_picture(img_path, left, top)     # Adding image to the blank slide

slide = prs.slides.add_slide(blank_slide_layout)        # Adding another slide
pic2= slide.shapes.add_picture(img_path2, left, top)    #Adding next image

prs.save('Presentation3.pptx')

