from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.chart.chart import Chart

# create presentation with 1 slide ------
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])

# define chart data ---------------------
chart_data = CategoryChartData()
chart_data.categories = ['East', 'West', 'Midwest']
chart_data.add_series('Series 1', (19,20,21))
# add chart to slide --------------------
x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
)

prs.save('chart-01.pptx')

Presentation('chart-01.pptx')
chart = prs.slides[0].shapes[1].chart
chart_data_old = chart.series[0].values

# Printing previous values
print(chart_data_old)

# Create new data to be plotted
chart_data = CategoryChartData()
chart_data.add_series('Series 2', (10, 20, 30))
Chart.replace_data(chart_data)
x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
)
prs.save('chart-02.pptx')